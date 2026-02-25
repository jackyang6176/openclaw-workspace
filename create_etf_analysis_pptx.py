from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 創建新的簡報
prs = Presentation()

# 設定slide寬高
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 標題頁
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "中國市場ETF投資組合分析"
subtitle.text = "00655L、00882、00887專業評估\n2026年2月"

# 目錄頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "目錄"
content = slide.placeholders[1]
content.text = "1. 持有ETF概覽\n2. 個別ETF深度分析\n3. 組合風險與機會評估\n4. 市場展望與操作建議\n5. 投資策略調整方向"

# ETF概覽頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "持有ETF概覽"
content = slide.placeholders[1]
content.text = "• 00655L（國泰A50正2）\n  - 槓桿型ETF，追蹤富時中國A50指數每日正向2倍報酬\n  - 高波動性，適合短期交易\n\n• 00882（中信中國高股息）\n  - 追蹤中證中國內地高股息精選指數\n  - 提供穩定股息收入，相對穩健\n\n• 00887（永豐中國科技50大）\n  - 追蹤中證海外中國科技50指數\n  - 投資海外上市中國科技龍頭企業"

# 00655L分析頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "00655L（國泰A50正2）深度分析"
content = slide.placeholders[1]
content.text = "• 產品特性：\n  - 每日重新平衡的2倍槓桿ETF\n  - 追蹤富時中國A50指數\n  - 成分股為中國A股市值最大的50家公司\n\n• 風險考量：\n  - 槓桿效應放大損益\n  - 長期持有可能因波動耗損產生偏離\n  - 適合短期交易，不適合長期投資\n\n• 操作建議：\n  - 建議密切監控市場動向\n  - 設定明確停損點\n  - 避免在高波動期間長期持有"

# 00882分析頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "00882（中信中國高股息）深度分析"
content = slide.placeholders[1]
content.text = "• 產品特性：\n  - 追蹤中證中國內地高股息精選指數\n  - 篩選標準：高股息率、財務穩健、流動性佳\n  - 成分股涵蓋金融、能源、公用事業等產業\n\n• 優勢分析：\n  - 提供穩定現金流\n  - 相對低波動性\n  - 在市場下跌時具防禦性\n\n• 操作建議：\n  - 可作為核心持股長期持有\n  - 關注中國經濟復甦進度\n  - 定期檢視成分股股息政策變化"

# 00887分析頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "00887（永豐中國科技50大）深度分析"
content = slide.placeholders[1]
content.text = "• 產品特性：\n  - 追蹤中證海外中國科技50指數\n  - 成分股包含阿里、騰訊、美團等科技巨頭\n  - 投資標的為海外上市中國科技企業\n\n• 市場機會：\n  - 中國科技產業長期成長潛力\n  - 數位轉型趨勢持續\n  - 監管環境逐漸明朗化\n\n• 風險考量：\n  - 地緣政治風險\n  - 中美關係影響\n  - 科技監管政策變化\n\n• 操作建議：\n  - 分批布局降低風險\n  - 關注企業基本面改善\n  - 配合整體科技產業景氣循環"

# 組合風險評估頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "投資組合風險與機會評估"
content = slide.placeholders[1]
content.text = "• 風險分析：\n  - 高度集中於中國市場\n  - 00655L槓桿風險\n  - 地緣政治不確定性\n  - 人民幣匯率波動\n\n• 機會評估：\n  - 中國經濟復甦潛力\n  - 科技產業創新動能\n  - 高股息提供穩定收益\n  - A50指數代表中國核心資產\n\n• 相關性分析：\n  - 三檔ETF高度相關（同屬中國市場）\n  - 建議增加地域分散性"

# 市場展望頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "中國市場展望與操作建議"
content = slide.placeholders[1]
content.text = "• 宏觀經濟展望：\n  - 中國經濟復甦力度待觀察\n  - 政策支持持續但效果有限\n  - 房地產市場調整尚未結束\n\n• 操作建議：\n  - 00655L：短期交易，嚴格控管風險\n  - 00882：可長期持有，作為收益來源\n  - 00887：分批布局，關注科技監管政策\n\n• 資金配置建議：\n  - 降低槓桿部位比例\n  - 增加非中國市場曝險\n  - 保持現金部位靈活性"

# 投資策略調整方向
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "投資策略調整方向"
content = slide.placeholders[1]
content.text = "• 短期策略（1-3個月）：\n  - 00655L設定明確停損點（建議-15%）\n  - 關注中國兩會政策方向\n  - 觀察美聯儲利率決策對新興市場影響\n\n• 中期策略（3-12個月）：\n  - 逐步降低槓桿部位至總投資10%以下\n  - 增加全球分散投資（如全球科技、美國市場）\n  - 保留00882作為核心收益部位\n\n• 長期策略（1年以上）：\n  - 建立多元地域投資組合\n  - 定期再平衡\n  - 關注ESG投資趨勢\n  - 考慮增加債券部位降低波動"

# 結論頁
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "結論與行動要點"
content = slide.placeholders[1]
content.text = "• 當前組合優點：\n  - 涵蓋中國市場多面向（大盤、高股息、科技）\n  - 00882提供穩定收益\n\n• 主要風險：\n  - 過度集中中國市場\n  - 00655L槓桿風險過高\n\n• 行動要點：\n  1. 立即檢視00655L部位並設定停損\n  2. 保留00882作為核心持股\n  3. 00887可逢低分批加碼\n  4. 規劃增加非中國市場投資\n\n• 定期檢視頻率：每月一次"

# 儲存簡報
prs.save('中國市場ETF投資組合分析.pptx')
print("簡報已成功創建：中國市場ETF投資組合分析.pptx")