#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_etf_presentation():
    # 創建新的簡報
    prs = Presentation()
    
    # 設定寬高 (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 第一頁：標題頁
    slide_layout = prs.slide_layouts[0]  # 標題版面
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "大陸ETF投資分析簡報"
    subtitle.text = "00655L、00882、00887 持有分析與建議\n2026年2月"
    
    # 第二頁：目錄
    slide_layout = prs.slide_layouts[1]  # 標題和內容版面
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "目錄"
    content.text = "• ETF基本資訊概覽\n• 各檔ETF詳細分析\n• 風險評估\n• 投資組合配置建議\n• 後續操作策略\n• 結論與展望"
    
    # 第三頁：ETF基本資訊概覽
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ETF基本資訊概覽"
    content.text = "• 00655L - 元大S&P中國A股低波動證券投資信託基金\n  - 追蹤指數：S&P中國A股低波動率指數\n  - 投資標的：中國A股市場中波動率較低的股票\n  - 特色：降低波動風險，適合穩健型投資人\n\n• 00882 - 富邦中國A50全收益指數基金\n  - 追蹤指數：富時中國A50全收益指數\n  - 投資標的：中國A股市值最大的50家公司\n  - 特色：涵蓋大型藍籌股，代表性強\n\n• 00887 - 國泰中國A50 ETF\n  - 追蹤指數：MSCI中國A50指數\n  - 投資標的：MSCI中國A50指數成分股\n  - 特色：國際標準指數，全球資金配置參考"
    
    # 第四頁：00655L詳細分析
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "00655L - 元大S&P中國A股低波動ETF分析"
    content.text = "• 投資策略：選取波動率最低的100檔中國A股\n• 成分股特色：金融、消費、公用事業等防禦型產業比重高\n• 費率結構：經理費0.3%，保管費0.035%\n• 流動性：日均成交量約XX萬單位\n• 近期表現：\n  - 年初至今報酬率：XX%\n  - 一年報酬率：XX%\n  - 波動率：XX% (低於大盤)\n• 適合投資人：追求穩定收益、風險承受度較低者"
    
    # 第五頁：00882詳細分析
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "00882 - 富邦中國A50全收益ETF分析"
    content.text = "• 投資策略：完全複製富時中國A50全收益指數\n• 成分股特色：大型金融、科技、消費龍頭企業\n• 費率結構：經理費0.315%，保管費0.035%\n• 流動性：日均成交量約XX萬單位\n• 近期表現：\n  - 年初至今報酬率：XX%\n  - 一年報酬率：XX%\n  - 股息收益率：XX%\n• 適合投資人：看好中國大型企業長期發展者"
    
    # 第六頁：00887詳細分析
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "00887 - 國泰中國A50 ETF分析"
    content.text = "• 投資策略：追蹤MSCI中國A50指數\n• 成分股特色：均衡配置各產業龍頭，包含新經濟企業\n• 費率結構：經理費0.3%，保管費0.035%\n• 流動性：日均成交量約XX萬單位\n• 近期表現：\n  - 年初至今報酬率：XX%\n  - 一年報酬率：XX%\n  - 追蹤誤差：XX%\n• 適合投資人：希望獲得中國市場整體成長機會者"
    
    # 第七頁：風險評估
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "風險評估"
    content.text = "• 政策風險：中國監管政策變化對市場影響大\n• 匯率風險：人民幣匯率波動影響報酬\n• 流動性風險：部分ETF成交量較低，買賣價差較大\n• 追蹤誤差：不同ETF追蹤效率有所差異\n• 市場風險：中國經濟放緩可能影響企業獲利\n• 地緣政治風險：中美關係緊張影響外資信心\n\n風險評級：中高風險 (相較於台灣或美國市場)"
    
    # 第八頁：投資組合配置建議
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "投資組合配置建議"
    content.text = "• 現有持倉分析：\n  - 三檔ETF皆聚焦中國A股，集中度過高\n  - 00882與00887成分股重疊度高，分散效果有限\n  - 00655L提供低波動特性，但整體仍偏高風險\n\n• 建議配置比例：\n  - 中國A股ETF總部位建議不超過投資組合15-20%\n  - 若已超標，建議適度減碼\n  - 可考慮保留00655L作為防禦配置，精簡其他兩檔\n\n• 替代選擇：\n  - 考慮增加港股或全球新興市場ETF分散風險"
    
    # 第九頁：後續操作策略
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "後續操作策略"
    content.text = "• 短期策略（1-3個月）：\n  - 觀察中國兩會政策方向\n  - 關注美中關係發展\n  - 評估人民幣匯率走勢\n\n• 中期策略（3-12個月）：\n  - 若中國經濟數據改善，可考慮加碼\n  - 若市場波動加劇，優先保留00655L\n  - 定期檢視持股比例，避免單一市場過度集中\n\n• 長期策略（1年以上）：\n  - 中國市場長期仍有成長潛力\n  - 建議採用定期定額方式累積部位\n  - 持續關注ESG趨勢對中國企業影響"
    
    # 第十頁：結論與展望
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "結論與展望"
    content.text = "• 投資價值：\n  - 中國A股估值處於歷史相對低位\n  - 政策支持有望帶動市場回升\n  - 長期而言仍具配置價值\n\n• 操作建議：\n  - 精簡持倉：保留00655L，評估是否需同時持有00882與00887\n  - 控制部位：中國A股ETF總部位建議控制在15-20%以內\n  - 分散風險：考慮增加其他區域或資產類別配置\n\n• 展望：\n  - 2026年中國經濟復甦力度將是關鍵\n  - 關注消費、科技、綠能等政策支持產業\n  - 長期投資人可逢低布局，短期注意波動風險"
    
    # 儲存簡報
    prs.save('大陸ETF投資分析簡報.pptx')
    print("簡報已成功建立：大陸ETF投資分析簡報.pptx")

if __name__ == "__main__":
    create_etf_presentation()