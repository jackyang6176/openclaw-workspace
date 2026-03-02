#!/usr/bin/env python3
"""
合約管理系統專案提案簡報生成器
使用 python-pptx 創建專業簡報
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 顏色定義
CORPORATE_BLUE = RGBColor(0, 120, 212)  # #0078D4
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)

def create_title_slide(prs, title, subtitle, date):
    """創建標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # 背景色塊
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CORPORATE_BLUE
    shape.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_section_slide(prs, section_num, section_title):
    """創建章節頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 左側色塊
    left = Inches(0)
    top = Inches(0)
    width = Inches(1.5)
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CORPORATE_BLUE
    shape.line.fill.background()
    
    # 章節編號
    num_box = slide.shapes.add_textbox(Inches(0.3), Inches(2.5), Inches(1), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 章節標題
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(7.5), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    return slide

def create_content_slide(prs, title, content_items):
    """創建內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CORPORATE_BLUE
    
    # 分隔線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = CORPORATE_BLUE
    line.line.fill.background()
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
    
    return slide

def create_table_slide(prs, title, headers, data):
    """創建表格頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CORPORATE_BLUE
    
    # 表格
    rows = len(data) + 1
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 設置列寬
    for i in range(cols):
        table.columns[i].width = Inches(9 / cols)
    
    # 填充標題列
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CORPORATE_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    
    # 填充數據列
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def create_timeline_slide(prs, title, phases):
    """創建時程圖頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CORPORATE_BLUE
    
    # 時間軸
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1)
    
    num_phases = len(phases)
    phase_width = width / num_phases
    
    for i, phase in enumerate(phases):
        phase_left = left + (i * phase_width)
        
        # 階段方塊
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            phase_left + Inches(0.05),
            top,
            phase_width - Inches(0.1),
            height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CORPORATE_BLUE
        shape.line.color.rgb = DARK_BLUE
        
        # 階段文字
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"第{phase['month']}月\n{phase['name']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_proposal():
    """生成完整簡報"""
    prs = Presentation()
    
    # 設置寬螢幕 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print("📊 正在生成簡報...")
    
    # 1. 標題頁
    print("  ✓ 標題頁")
    create_title_slide(
        prs,
        "合約管理系統專案提案",
        "智能合約生命週期管理解決方案",
        "2026年3月"
    )
    
    # 2. 章節1：公司背景
    print("  ✓ 章節1：公司背景")
    create_section_slide(prs, "01", "公司背景與相關專案經驗")
    
    create_content_slide(prs, "核心優勢", [
        "多年企業級系統開發經驗",
        "完整的合約管理領域知識",
        "熟悉 M365、AD、SSO 整合技術",
        "專業的 AI 文件處理能力"
    ])
    
    create_content_slide(prs, "相關專案經驗", [
        "大型企業電子簽核系統建置",
        "文件管理與版本控制系統",
        "HR 系統整合與單一登入 (SSO) 專案",
        "AI 輔助文件比對與 OCR 應用"
    ])
    
    # 3. 章節2：需求理解
    print("  ✓ 章節2：需求理解")
    create_section_slide(prs, "02", "對本專案需求之理解")
    
    create_content_slide(prs, "兩大核心系統", [
        "線上合約建立系統",
        "合約管理中心",
        "權限分離管理（流程權限 vs 歸檔後權限）",
        "自動稽催機制",
        "AI 文件比對功能",
        "多系統整合（HR、ISIM、影印機）"
    ])
    
    # 4. 章節3：系統功能
    print("  ✓ 章節3：系統功能")
    create_section_slide(prs, "03", "系統功能與架構說明")
    
    # 功能表格
    headers = ["模組", "功能項目", "工期"]
    data = [
        ["合約建立與簽核", "合約審核管理表單", "1月"],
        ["合約建立與簽核", "制式合約表單", "1月"],
        ["合約建立與簽核", "用印申請管理", "0.5月"],
        ["AI 智能應用", "合約文件比對", "2月"],
        ["通知管理", "提醒與通知機制", "0.5月"],
        ["報表中心", "合約相關報表", "0.5月"],
        ["系統整合", "HR/ISIM/API整合", "2月"],
        ["參數配置", "表單參數設計", "0.5月"],
    ]
    create_table_slide(prs, "完整功能模組（總工期8個月）", headers, data)
    
    # 5. 章節4：操作流程
    print("  ✓ 章節4：操作流程")
    create_section_slide(prs, "04", "系統操作流程與使用情境")
    
    create_content_slide(prs, "典型使用情境", [
        "情境A：新合約建立流程（承辦人→主管→法務→高階→用印→簽約→歸檔）",
        "情境B：合約到期管理（90天前提醒→30天前再次提醒→續約確認）",
        "情境C：AI 文件比對（電子檔 + 紙本掃描→AI 判讀→差異標註）"
    ])
    
    # 6. 章節5：專案管理
    print("  ✓ 章節5：專案管理")
    create_section_slide(prs, "05", "專案管理方式與執行流程")
    
    # 時程圖
    phases = [
        {"month": "1-2", "name": "需求確認\n系統設計"},
        {"month": "3-4", "name": "核心功能\n開發"},
        {"month": "5-6", "name": "AI 比對\n系統整合"},
        {"month": "7-8", "name": "測試驗收\n正式上線"},
    ]
    create_timeline_slide(prs, "開發階段規劃（8個月）", phases)
    
    create_content_slide(prs, "交付里程碑", [
        "Phase 1（第1-2月）：需求確認、系統設計、UI/UX 原型",
        "Phase 2（第3-4月）：核心功能開發（合約管理、簽核流程）",
        "Phase 3（第5-6月）：AI 比對、系統整合、報表功能",
        "Phase 4（第7-8月）：測試驗收、教育訓練、正式上線"
    ])
    
    # 7. 章節6：資訊安全
    print("  ✓ 章節6：資訊安全")
    create_section_slide(prs, "06", "資訊安全、權限控管及法遵")
    
    create_content_slide(prs, "多層次安全防護", [
        "ISIM 整合：統一密碼管理，符合企業資安政策",
        "SSO 單一登入：無縫整合 M365、AD",
        "多因子驗證：強化帳戶安全",
        "浮水印保護機密文件",
        "完整稽核軌跡記錄",
        "版本控制與修改歷程追蹤"
    ])
    
    # 8. 章節7：導入方式
    print("  ✓ 章節7：導入方式")
    create_section_slide(prs, "07", "導入方式、時程與資源配置")
    
    create_content_slide(prs, "漸進式上線計畫", [
        "第7月：試營運（UAT、種子用戶測試、問題修復）",
        "第8月：正式上線（全面啟用、教育訓練、平行運作）"
    ])
    
    create_content_slide(prs, "資源配置建議", [
        "專案經理：1人（全程8個月）",
        "系統分析師：2人（第1-4月）",
        "前端工程師：2人（第2-7月）",
        "後端工程師：3人（第2-7月）",
        "AI 工程師：1人（第3-5月）",
        "測試工程師：2人（第5-8月）"
    ])
    
    # 9. 結語
    print("  ✓ 結語")
    create_content_slide(prs, "為什麼選擇我們？", [
        "🎯 完整解決方案：從合約建立到歸檔的全生命週期管理",
        "🚀 創新 AI 應用：智能文件比對，提升審查效率",
        "🔒 企業級安全：多層次權限控管，符合資安規範",
        "🔗 無縫整合：與現有 HR、ISIM 系統完美整合"
    ])
    
    # 10. 感謝頁
    print("  ✓ 感謝頁")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色塊
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CORPORATE_BLUE
    shape.line.fill.background()
    
    # 感謝文字
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感謝聆聽"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "期待與您攜手打造智能合約管理新時代"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 儲存檔案
    output_path = "/home/admin/.openclaw/workspace/contract-proposal/proposal.pptx"
    prs.save(output_path)
    print(f"\n✅ 簡報已生成：{output_path}")
    print(f"📊 總頁數：{len(prs.slides)} 頁")
    
    return output_path

if __name__ == "__main__":
    generate_proposal()
